"""Engine-level tests for immersive reading: extract → store → search → export.

These exercise the pure engine against a temp root, so nothing here needs the
path service, a user workspace, or an LLM.
"""

from __future__ import annotations

import json
from pathlib import Path
import zipfile

import pytest

from deeptutor.reading import (
    Annotation,
    MaterialNotFound,
    ReadingError,
    ReadingStore,
    Rect,
    TextPositionSelector,
    TextQuoteSelector,
    export_material,
    parse_locators,
    render_outline,
    render_units,
    search_material,
    verify_quote,
)
from deeptutor.reading.extract import (
    SECTION_TARGET_CHARS,
    extract_material,
    first_line_label,
    split_into_sections,
)
from deeptutor.reading.models import parse_text_selectors
from deeptutor.reading.search import normalise, search_units, terms_of

pymupdf = pytest.importorskip("pymupdf")


# ---------------------------------------------------------------------------
# fixtures
# ---------------------------------------------------------------------------


def _write_pdf(path: Path, pages: list[str], *, toc: list | None = None) -> Path:
    doc = pymupdf.open()
    for body in pages:
        page = doc.new_page()
        page.insert_textbox(pymupdf.Rect(50, 50, 545, 780), body, fontsize=11)
    if toc:
        doc.set_toc(toc)
    doc.save(path)
    doc.close()
    return path


def _write_epub(path: Path) -> Path:
    with zipfile.ZipFile(path, "w") as archive:
        archive.writestr("mimetype", "application/epub+zip")
        archive.writestr(
            "META-INF/container.xml",
            """<?xml version="1.0"?><container xmlns="urn:oasis:names:tc:opendocument:xmlns:container"><rootfiles><rootfile full-path="OEBPS/content.opf" media-type="application/oebps-package+xml"/></rootfiles></container>""",
        )
        archive.writestr(
            "OEBPS/content.opf",
            """<?xml version="1.0"?><package xmlns="http://www.idpf.org/2007/opf" version="3.0"><manifest><item id="nav" href="nav.xhtml" media-type="application/xhtml+xml" properties="nav"/><item id="one" href="chapters/one.xhtml" media-type="application/xhtml+xml"/><item id="two" href="chapters/two.xhtml" media-type="application/xhtml+xml"/></manifest><spine><itemref idref="one"/><itemref idref="two"/></spine></package>""",
        )
        archive.writestr(
            "OEBPS/nav.xhtml",
            """<html xmlns="http://www.w3.org/1999/xhtml" xmlns:epub="http://www.idpf.org/2007/ops"><body><nav epub:type="toc"><ol><li><a href="chapters/one.xhtml">Part One</a><ol><li><a href="chapters/two.xhtml">Second Chapter</a></li></ol></li></ol></nav></body></html>""",
        )
        archive.writestr(
            "OEBPS/chapters/one.xhtml",
            "<html xmlns='http://www.w3.org/1999/xhtml'><head><title>One</title></head><body><h1>First Chapter</h1><p>Alpha source text.</p><script>ignore me</script></body></html>",
        )
        archive.writestr(
            "OEBPS/chapters/two.xhtml",
            "<html xmlns='http://www.w3.org/1999/xhtml'><body><h1>Second Chapter</h1><p>Beta source text.</p></body></html>",
        )
    return path


@pytest.fixture
def pdf_path(tmp_path: Path) -> Path:
    return _write_pdf(
        tmp_path / "attention.pdf",
        [
            "Chapter one. Introduction to sequence models and their limits.",
            "Chapter two. Transformers use scaled dot-product attention.",
            "Chapter three. Positional encoding injects order information.",
        ],
        toc=[[1, "Introduction", 1], [1, "Transformers", 2], [2, "Positional encoding", 3]],
    )


@pytest.fixture
def store(tmp_path: Path) -> ReadingStore:
    return ReadingStore(root=tmp_path / "materials")


# ---------------------------------------------------------------------------
# extract
# ---------------------------------------------------------------------------


def test_pdf_extracts_one_unit_per_page_with_its_own_outline(pdf_path: Path) -> None:
    extraction = extract_material(pdf_path)

    assert extraction.unit == "page"
    assert len(extraction.units) == 3
    assert "scaled dot-product" in extraction.units[1]
    assert extraction.has_raw_view is True
    # The document's own bookmarks win over synthesised labels.
    assert [(e.locator, e.title) for e in extraction.outline] == [
        (1, "Introduction"),
        (2, "Transformers"),
        (3, "Positional encoding"),
    ]
    assert extraction.outline[2].level == 2


def test_pdf_outline_drops_bookmarks_pointing_outside_the_page_range(tmp_path: Path) -> None:
    path = _write_pdf(tmp_path / "bad_toc.pdf", ["only page"], toc=[[1, "Ghost", 1]])
    doc = pymupdf.open(path)
    # Rewrite the bookmark to a page that does not exist, the way some
    # generators do; the extractor must skip it rather than clamp it.
    doc.set_toc([[1, "Ghost", 1]])
    doc.save(tmp_path / "bad_toc2.pdf")
    doc.close()

    extraction = extract_material(tmp_path / "bad_toc2.pdf")
    assert all(1 <= e.locator <= len(extraction.units) for e in extraction.outline)


def test_pdf_without_bookmarks_does_not_invent_contents(tmp_path: Path) -> None:
    path = _write_pdf(
        tmp_path / "plain.pdf",
        ["Figure 1: this is page content, not a document heading.", "References"],
    )

    extraction = extract_material(path)
    assert extraction.outline == ()

    reading_store = ReadingStore(root=tmp_path / "plain-materials")
    manifest = reading_store.ingest(path)
    assert reading_store.outline(manifest.material_id) == []


def test_text_file_is_cut_into_sections_on_paragraph_boundaries(tmp_path: Path) -> None:
    paragraph = "Dense prose about attention mechanisms. " * 30  # ~1.2k chars
    path = tmp_path / "notes.md"
    path.write_text("\n\n".join([paragraph] * 8), encoding="utf-8")

    extraction = extract_material(path)

    assert extraction.unit == "section"
    assert extraction.has_raw_view is False
    assert len(extraction.units) > 1
    # Cuts land on paragraph boundaries, so no unit starts mid-sentence.
    assert all(unit.startswith("Dense prose") for unit in extraction.units)


def test_epub_preserves_spine_units_source_hrefs_and_nested_outline(tmp_path: Path) -> None:
    extraction = extract_material(_write_epub(tmp_path / "book.epub"))

    assert extraction.render_mode == "epub"
    assert extraction.has_raw_view is False
    assert extraction.unit == "chapter"
    assert extraction.units == (
        "First Chapter\nAlpha source text.",
        "Second Chapter\nBeta source text.",
    )
    assert [ref.source_href for ref in extraction.unit_refs] == [
        "OEBPS/chapters/one.xhtml",
        "OEBPS/chapters/two.xhtml",
    ]
    assert [(row.locator, row.title, row.level) for row in extraction.outline] == [
        (1, "Part One", 1),
        (2, "Second Chapter", 2),
    ]


def test_pptx_slides_become_units_when_the_extractor_marks_them(tmp_path: Path) -> None:
    pytest.importorskip("pptx")
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    for text in ("First slide body", "Second slide body"):
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
        box.text_frame.text = text
    path = tmp_path / "deck.pptx"
    prs.save(path)

    extraction = extract_material(path)

    assert extraction.unit == "slide"
    assert len(extraction.units) == 2
    assert "First slide body" in extraction.units[0]


def test_empty_and_unreadable_sources_raise_reading_error(tmp_path: Path) -> None:
    missing = tmp_path / "nope.pdf"
    with pytest.raises(ReadingError):
        extract_material(missing)

    blank = tmp_path / "blank.txt"
    blank.write_text("   \n\n  ", encoding="utf-8")
    with pytest.raises(ReadingError):
        extract_material(blank)


def test_a_single_enormous_line_is_still_split(tmp_path: Path) -> None:
    path = tmp_path / "minified.txt"
    path.write_text("x" * (SECTION_TARGET_CHARS * 4), encoding="utf-8")

    units = split_into_sections(path.read_text(encoding="utf-8"))

    assert len(units) > 1
    assert all(units)


def test_first_line_label_prefers_a_markdown_heading() -> None:
    assert first_line_label("## Attention\n\nbody text") == "Attention"
    assert first_line_label("plain first line\nsecond") == "plain first line"
    assert first_line_label("\n\n") == ""


# ---------------------------------------------------------------------------
# store
# ---------------------------------------------------------------------------


def test_ingest_writes_units_raw_and_manifest(store: ReadingStore, pdf_path: Path) -> None:
    manifest = store.ingest(pdf_path)

    assert manifest.unit_count == 3
    assert manifest.has_raw_view is True
    assert manifest.filename == "attention.pdf"
    assert store.exists(manifest.material_id)
    assert store.unit_text(manifest.material_id, 2).find("scaled dot-product") >= 0
    raw = store.raw_path(manifest.material_id)
    assert raw is not None and raw.read_bytes()[:5] == b"%PDF-"


def test_epub_store_keeps_original_but_legacy_pdf_flag_stays_false(
    store: ReadingStore, tmp_path: Path
) -> None:
    path = _write_epub(tmp_path / "book.epub")
    manifest = store.ingest(path)

    assert manifest.render_mode == "epub"
    assert manifest.has_raw_view is False
    assert store.raw_path(manifest.material_id) is not None
    assert store.unit_references(manifest.material_id)[1].source_href.endswith("two.xhtml")


def test_position_round_trip_validates_locator_and_anchor(
    store: ReadingStore, tmp_path: Path
) -> None:
    from deeptutor.reading import ReadingPosition

    manifest = store.ingest(_write_epub(tmp_path / "book.epub"))
    saved = store.save_position(
        manifest.material_id,
        ReadingPosition(locator=2, source_anchor="epubcfi(/6/4)", percentage=0.6),
    )

    assert saved.updated_at > 0
    assert store.position(manifest.material_id).source_anchor == "epubcfi(/6/4)"
    with pytest.raises(ReadingError):
        store.save_position(manifest.material_id, ReadingPosition(locator=99))


def test_bookmarks_are_plural_deliberate_and_deduplicated_by_locator(
    store: ReadingStore, tmp_path: Path
) -> None:
    """Unlike the position, these are chosen — several of them, by id."""
    manifest = store.ingest(_write_epub(tmp_path / "book.epub"))
    material_id = manifest.material_id

    assert store.bookmarks(material_id) == []

    second = store.add_bookmark(material_id, 2, "the good bit")
    first = store.add_bookmark(material_id, 1)

    assert second.label == "the good bit"
    # Saving a place must not force the reader to name it first.
    assert first.label == ""
    assert first.bookmark_id != second.bookmark_id

    # Reading order, not insertion order: the list is an index, not a log.
    assert [row.locator for row in store.bookmarks(material_id)] == [1, 2]

    # The affordance is a toggle on the reader's toolbar, so "bookmark this
    # page" when the page is already bookmarked is the existing bookmark —
    # not a second identical row.
    again = store.add_bookmark(material_id, 2, "ignored")
    assert again.bookmark_id == second.bookmark_id
    assert again.label == "the good bit"
    assert len(store.bookmarks(material_id)) == 2

    with pytest.raises(ReadingError):
        store.add_bookmark(material_id, 99)

    assert store.delete_bookmark(material_id, second.bookmark_id) is True
    assert store.delete_bookmark(material_id, second.bookmark_id) is False
    assert [row.locator for row in store.bookmarks(material_id)] == [1]


def test_bookmarks_survive_a_compatible_reingest(store: ReadingStore, tmp_path: Path) -> None:
    """A bookmark is a place the reader chose, so a repair keeps it.

    The automatic viewport is allowed to reset when the spine changes under
    it; a deliberately-kept place is user-owned state like an annotation.
    """
    source = _write_epub(tmp_path / "book.epub")
    manifest = store.ingest(source)
    kept = store.add_bookmark(manifest.material_id, 2, "chapter two")

    reingested = store.ingest(source)

    rows = store.bookmarks(reingested.material_id)
    assert [row.bookmark_id for row in rows] == [kept.bookmark_id]
    assert rows[0].label == "chapter two"


def _downgrade_epub_manifest_to_legacy_text(store: ReadingStore, material_id: str) -> None:
    material_dir = store.root / material_id
    manifest_path = material_dir / "manifest.json"
    data = json.loads(manifest_path.read_text(encoding="utf-8"))
    data.pop("render_mode", None)
    data["has_raw_view"] = False
    manifest_path.write_text(json.dumps(data), encoding="utf-8")
    raw_dir = material_dir / "raw"
    if raw_dir.exists():
        for child in raw_dir.iterdir():
            child.unlink()
        raw_dir.rmdir()


def test_reupload_upgrades_legacy_epub_when_it_has_no_annotations(
    store: ReadingStore, tmp_path: Path
) -> None:
    path = _write_epub(tmp_path / "book.epub")
    first = store.ingest(path)
    _downgrade_epub_manifest_to_legacy_text(store, first.material_id)

    upgraded = store.ingest(path)

    assert upgraded.render_mode == "epub"
    assert store.raw_path(upgraded.material_id) is not None


def test_legacy_epub_upgrade_discards_out_of_range_position(
    store: ReadingStore, tmp_path: Path
) -> None:
    from deeptutor.reading import ReadingPosition

    path = _write_epub(tmp_path / "book.epub")
    first = store.ingest(path)
    # Emulate a legacy text extraction with more units than the EPUB spine.
    material_dir = store.root / first.material_id
    manifest_path = material_dir / "manifest.json"
    data = json.loads(manifest_path.read_text(encoding="utf-8"))
    data["unit_count"] = 3
    manifest_path.write_text(json.dumps(data), encoding="utf-8")
    (material_dir / "units" / "0003.txt").write_text("Legacy split", encoding="utf-8")
    store.save_position(
        first.material_id, ReadingPosition(locator=3, source_anchor="old-text-anchor")
    )
    _downgrade_epub_manifest_to_legacy_text(store, first.material_id)

    upgraded = store.ingest(path)

    assert upgraded.unit_count == 2
    assert store.position(upgraded.material_id).locator == 1
    assert store.position(upgraded.material_id).source_anchor == ""


def test_reupload_rejects_legacy_epub_upgrade_with_annotations(
    store: ReadingStore, tmp_path: Path
) -> None:
    from deeptutor.reading import ReadingUpgradeConflict

    path = _write_epub(tmp_path / "book.epub")
    first = store.ingest(path)
    store.save_annotation(first.material_id, Annotation(annotation_id="", locator=1, quote="Alpha"))
    _downgrade_epub_manifest_to_legacy_text(store, first.material_id)

    with pytest.raises(ReadingUpgradeConflict):
        store.ingest(path)

    assert [row.quote for row in store.annotations(first.material_id)] == ["Alpha"]


def test_reingesting_the_same_bytes_reuses_the_material_and_its_annotations(
    store: ReadingStore, pdf_path: Path, tmp_path: Path
) -> None:
    first = store.ingest(pdf_path)
    store.save_annotation(
        first.material_id,
        Annotation(annotation_id="", locator=2, quote="attention", note="key idea"),
    )

    copy = tmp_path / "renamed.pdf"
    copy.write_bytes(pdf_path.read_bytes())
    second = store.ingest(copy)

    assert second.material_id == first.material_id
    assert [a.note for a in store.annotations(second.material_id)] == ["key idea"]


def test_unknown_material_and_bad_id_are_distinguishable(store: ReadingStore) -> None:
    with pytest.raises(MaterialNotFound):
        store.manifest("0123456789abcdef")
    with pytest.raises(ReadingError):
        store.manifest("../../etc/passwd")
    with pytest.raises(ReadingError):
        store.manifest("NOT-HEX")


def test_out_of_range_locator_reports_the_real_range(store: ReadingStore, pdf_path: Path) -> None:
    manifest = store.ingest(pdf_path)
    with pytest.raises(ReadingError) as excinfo:
        store.unit_text(manifest.material_id, 99)
    assert "3" in str(excinfo.value)


def test_read_units_is_bounded_and_says_so(store: ReadingStore, pdf_path: Path) -> None:
    manifest = store.ingest(pdf_path)

    rows, truncated = store.read_units(manifest.material_id, [1, 2, 3], max_chars=40)

    assert truncated is True
    assert sum(len(text) for _, text in rows) <= 40


def test_list_materials_is_newest_first_and_skips_junk(
    store: ReadingStore, pdf_path: Path, tmp_path: Path
) -> None:
    first = store.ingest(pdf_path)
    other = _write_pdf(tmp_path / "second.pdf", ["another document body"])
    second = store.ingest(other)
    (store.root / "not-a-material").mkdir(parents=True, exist_ok=True)

    ids = [m.material_id for m in store.list_materials()]

    assert set(ids) == {first.material_id, second.material_id}
    assert len(ids) == 2


def test_delete_removes_everything(store: ReadingStore, pdf_path: Path) -> None:
    manifest = store.ingest(pdf_path)
    assert store.delete(manifest.material_id) is True
    assert store.exists(manifest.material_id) is False
    assert store.delete(manifest.material_id) is False


def test_partial_ingest_is_repaired_on_the_next_upload(store: ReadingStore, pdf_path: Path) -> None:
    manifest = store.ingest(pdf_path)
    # Simulate a crash between unit writes: the last unit is gone but the
    # manifest still claims it.
    (store.root / manifest.material_id / "units" / "0003.txt").unlink()

    reingested = store.ingest(pdf_path)

    assert reingested.unit_count == 3
    assert store.unit_text(reingested.material_id, 3).strip() != ""


# ---------------------------------------------------------------------------
# annotations
# ---------------------------------------------------------------------------


def test_annotations_round_trip_with_generated_ids(store: ReadingStore, pdf_path: Path) -> None:
    manifest = store.ingest(pdf_path)

    saved = store.save_annotation(
        manifest.material_id,
        Annotation(
            annotation_id="",
            locator=1,
            quote="Introduction",
            note="start here",
            rects=(Rect(0.1, 0.2, 0.5, 0.25),),
        ),
    )

    assert saved.annotation_id
    assert saved.material_revision == manifest.revision
    stored = store.annotations(manifest.material_id)
    assert len(stored) == 1
    assert stored[0].rects[0].to_list() == [0.1, 0.2, 0.5, 0.25]


def test_w3c_text_selectors_round_trip_and_can_supply_the_quote(
    store: ReadingStore, pdf_path: Path
) -> None:
    manifest = store.ingest(pdf_path)
    unit_text = store.unit_text(manifest.material_id, 1)
    start = unit_text.index("Introduction")

    saved = store.save_annotation(
        manifest.material_id,
        Annotation(
            annotation_id="",
            locator=1,
            selectors=(
                TextQuoteSelector(exact="Introduction", suffix=" to sequence"),
                TextPositionSelector(start=start, end=start + 12),
            ),
        ),
    )

    assert saved.quote == "Introduction"
    assert [
        selector.to_dict() for selector in store.annotations(manifest.material_id)[0].selectors
    ] == [
        {
            "type": "TextQuoteSelector",
            "exact": "Introduction",
            "suffix": " to sequence",
        },
        {"type": "TextPositionSelector", "start": start, "end": start + 12},
    ]


def test_mismatched_quote_and_text_quote_selector_are_rejected(
    store: ReadingStore, pdf_path: Path
) -> None:
    manifest = store.ingest(pdf_path)

    with pytest.raises(ReadingError, match="does not match"):
        store.save_annotation(
            manifest.material_id,
            Annotation(
                annotation_id="",
                locator=1,
                quote="Introduction",
                selectors=(TextQuoteSelector(exact="different text"),),
            ),
        )


def test_text_quote_selector_must_occur_in_stored_unit(store: ReadingStore, pdf_path: Path) -> None:
    manifest = store.ingest(pdf_path)

    with pytest.raises(ReadingError, match="does not occur"):
        store.save_annotation(
            manifest.material_id,
            Annotation(
                annotation_id="",
                locator=1,
                selectors=(TextQuoteSelector(exact="not in this unit"),),
            ),
        )


def test_text_position_selector_cannot_extend_past_stored_unit(
    store: ReadingStore, pdf_path: Path
) -> None:
    manifest = store.ingest(pdf_path)
    unit_length = len(store.unit_text(manifest.material_id, 1))

    with pytest.raises(ReadingError, match="extends past"):
        store.save_annotation(
            manifest.material_id,
            Annotation(
                annotation_id="",
                locator=1,
                selectors=(TextPositionSelector(start=0, end=unit_length + 1),),
            ),
        )


def test_legacy_selector_parser_keeps_prefix_tail_nearest_the_quote() -> None:
    parsed = parse_text_selectors(
        [
            {
                "type": "TextQuoteSelector",
                "exact": "text",
                "prefix": "a" * 200 + "b" * 200,
            }
        ]
    )

    assert isinstance(parsed[0], TextQuoteSelector)
    assert parsed[0].prefix == "b" * 128


def test_selector_whitespace_is_canonicalised_to_stored_unit(
    store: ReadingStore, pdf_path: Path
) -> None:
    manifest = store.ingest(pdf_path)
    unit_path = store.root / manifest.material_id / "units" / "0001.txt"
    unit_path.write_text("Before Sequence\n\nmodels after", encoding="utf-8")

    saved = store.save_annotation(
        manifest.material_id,
        Annotation(
            annotation_id="",
            locator=1,
            quote="Sequence models",
            selectors=(TextQuoteSelector(exact="Sequence models", prefix="Before"),),
        ),
    )

    assert saved.quote == "Sequence\n\nmodels"
    assert saved.selectors[0].to_dict()["exact"] == "Sequence\n\nmodels"


def test_quote_and_position_selectors_must_identify_the_same_text(
    store: ReadingStore, pdf_path: Path
) -> None:
    manifest = store.ingest(pdf_path)

    with pytest.raises(ReadingError, match="different text"):
        store.save_annotation(
            manifest.material_id,
            Annotation(
                annotation_id="",
                locator=1,
                selectors=(
                    TextQuoteSelector(exact="Introduction"),
                    TextPositionSelector(start=0, end=7),
                ),
            ),
        )


def test_saving_the_same_id_updates_in_place_and_keeps_created_at(
    store: ReadingStore, pdf_path: Path
) -> None:
    manifest = store.ingest(pdf_path)
    first = store.save_annotation(
        manifest.material_id, Annotation(annotation_id="", locator=1, note="v1")
    )

    updated = store.save_annotation(
        manifest.material_id,
        Annotation(annotation_id=first.annotation_id, locator=1, note="v2"),
    )

    assert updated.created_at == first.created_at
    assert updated.updated_at >= first.updated_at
    assert [a.note for a in store.annotations(manifest.material_id)] == ["v2"]


def test_annotation_on_a_nonexistent_locator_is_rejected(
    store: ReadingStore, pdf_path: Path
) -> None:
    manifest = store.ingest(pdf_path)
    with pytest.raises(ReadingError):
        store.save_annotation(manifest.material_id, Annotation(annotation_id="", locator=42))


def test_delete_annotation_reports_whether_it_existed(store: ReadingStore, pdf_path: Path) -> None:
    manifest = store.ingest(pdf_path)
    saved = store.save_annotation(
        manifest.material_id, Annotation(annotation_id="", locator=1, note="x")
    )

    assert store.delete_annotation(manifest.material_id, saved.annotation_id) is True
    assert store.delete_annotation(manifest.material_id, saved.annotation_id) is False
    assert store.annotations(manifest.material_id) == []


def test_malformed_rects_and_colors_are_normalised_not_trusted() -> None:
    parsed = Annotation.from_dict(
        {
            "annotation_id": "a1",
            "locator": 2,
            "kind": "scribble",
            "color": "neon",
            "rects": [[0.9, 0.9, 0.1, 0.1], "garbage", [0, 0, 0, 0], {"x0": 0, "y0": 0}],
        }
    )

    assert parsed.kind == "highlight"
    assert parsed.color == "yellow"
    # Inverted rect is ordered; degenerate and unparseable rects are dropped.
    assert [r.to_list() for r in parsed.rects] == [[0.1, 0.1, 0.9, 0.9]]


# ---------------------------------------------------------------------------
# search
# ---------------------------------------------------------------------------


def test_exact_search_returns_the_locator_and_a_snippet(
    store: ReadingStore, pdf_path: Path
) -> None:
    manifest = store.ingest(pdf_path)

    result = search_material(store, manifest.material_id, "scaled dot-product")

    assert result.mode == "exact"
    assert [hit.locator for hit in result.hits] == [2]
    assert "dot-product" in result.hits[0].snippet


def test_search_tolerates_line_wrapped_quotes() -> None:
    units = [(1, "Transformers use scaled\ndot-product   attention today.")]

    result = search_units(units, "scaled dot-product attention")

    assert result.mode == "normalised"
    assert result.hits[0].locator == 1


def test_term_ranking_is_the_fallback_and_prefers_more_matches() -> None:
    units = [
        (1, "positional encoding only"),
        (2, "attention and positional encoding together"),
        (3, "nothing relevant here"),
    ]

    result = search_units(units, "attention positional encoding")

    assert result.mode == "terms"
    assert result.hits[0].locator == 2


def test_search_returns_empty_for_a_blank_or_unmatched_query() -> None:
    units = [(1, "alpha beta")]
    assert search_units(units, "   ").is_empty
    assert search_units(units, "zzzzqqq").is_empty


def test_cjk_queries_are_bigram_expanded_so_they_match_partially() -> None:
    assert "注意" in terms_of("注意力机制")
    units = [(1, "本页讨论注意力机制的实现"), (2, "无关内容")]

    result = search_units(units, "注意力机制的推导")

    assert result.hits[0].locator == 1


def test_normalise_softens_quotes_and_whitespace() -> None:
    assert normalise("“Hello,  world!”") == normalise("Hello world")


def test_search_marks_truncation_when_more_hits_exist() -> None:
    units = [(i, "needle here") for i in range(1, 8)]

    result = search_units(units, "needle", limit=3)

    assert len(result.hits) == 3
    assert result.truncated is True


# ---------------------------------------------------------------------------
# service
# ---------------------------------------------------------------------------


@pytest.mark.parametrize(
    ("spec", "expected"),
    [
        ("2", [2]),
        (2, [2]),
        ("1-3", [1, 2, 3]),
        ("1–3", [1, 2, 3]),
        ("3,1", [1, 3]),
        ("1, 1, 2", [1, 2]),
        ([3, 2], [2, 3]),
        ("2-1", [1, 2]),
    ],
)
def test_parse_locators_accepts_the_grammar_the_model_types(spec, expected) -> None:
    assert parse_locators(spec, unit_count=3) == expected


def test_parse_locators_drops_out_of_range_and_raises_when_nothing_is_left() -> None:
    assert parse_locators("2,99", unit_count=3) == [2]
    with pytest.raises(ReadingError):
        parse_locators("99", unit_count=3)
    with pytest.raises(ReadingError):
        parse_locators("garbage", unit_count=3)


def test_parse_locators_bounds_an_absurd_range_without_materialising_it() -> None:
    assert len(parse_locators("1-100000", unit_count=500)) <= 24


def test_render_units_labels_by_unit_kind(store: ReadingStore, pdf_path: Path) -> None:
    manifest = store.ingest(pdf_path)

    rendered = render_units(store, manifest.material_id, "1-2")

    assert "--- Page 1 ---" in rendered.text
    assert "--- Page 2 ---" in rendered.text
    assert rendered.locators == (1, 2)
    assert rendered.truncated is False


def test_render_units_announces_truncation(store: ReadingStore, pdf_path: Path) -> None:
    manifest = store.ingest(pdf_path)

    rendered = render_units(store, manifest.material_id, "1-3", max_chars=30)

    assert rendered.truncated is True
    assert "truncated" in rendered.text


def test_render_outline_uses_the_documents_own_titles(store: ReadingStore, pdf_path: Path) -> None:
    manifest = store.ingest(pdf_path)

    text = render_outline(store, manifest.material_id)

    assert "page 2: Transformers" in text
    assert "attention.pdf" in text


def test_render_outline_falls_back_to_first_lines(store: ReadingStore, tmp_path: Path) -> None:
    path = tmp_path / "plain.md"
    path.write_text("# Alpha\nbody\n\n" + ("filler. " * 500) + "\n\n# Beta\nmore", encoding="utf-8")
    manifest = store.ingest(path)

    text = render_outline(store, manifest.material_id)

    assert "Alpha" in text


def test_verify_quote_confirms_a_real_quote(store: ReadingStore, pdf_path: Path) -> None:
    manifest = store.ingest(pdf_path)

    check = verify_quote(store, manifest.material_id, 2, "scaled dot-product")

    assert check.verified is True
    assert check.moved is False


def test_verify_quote_finds_the_right_locator_when_the_model_guessed_wrong(
    store: ReadingStore, pdf_path: Path
) -> None:
    manifest = store.ingest(pdf_path)

    check = verify_quote(store, manifest.material_id, 1, "scaled dot-product")

    assert check.verified is True
    assert check.found_locator == 2
    assert check.moved is True


def test_verify_quote_rejects_a_hallucinated_quote(store: ReadingStore, pdf_path: Path) -> None:
    manifest = store.ingest(pdf_path)

    check = verify_quote(store, manifest.material_id, 2, "quantum flux capacitor")

    assert check.verified is False


# ---------------------------------------------------------------------------
# export
# ---------------------------------------------------------------------------


def test_pdf_export_writes_real_annotations_back_into_the_file(
    store: ReadingStore, pdf_path: Path
) -> None:
    manifest = store.ingest(pdf_path)
    store.save_annotation(
        manifest.material_id,
        Annotation(
            annotation_id="",
            locator=2,
            kind="highlight",
            color="green",
            quote="scaled dot-product",
            note="the core mechanism",
            rects=(Rect(0.1, 0.1, 0.8, 0.15),),
        ),
    )

    result = export_material(store, manifest.material_id, fmt="pdf")

    assert result.filename == "attention-annotated.pdf"
    assert result.media_type == "application/pdf"
    with pymupdf.open(stream=result.data, filetype="pdf") as doc:
        annots = list(doc[1].annots())
        assert len(annots) == 1
        assert annots[0].info.get("content") == "the core mechanism"


def test_pdf_export_with_no_annotations_returns_the_original_bytes(
    store: ReadingStore, pdf_path: Path
) -> None:
    manifest = store.ingest(pdf_path)

    result = export_material(store, manifest.material_id, fmt="pdf")

    assert result.data == pdf_path.read_bytes()


def test_pdf_export_survives_one_unusable_annotation(store: ReadingStore, pdf_path: Path) -> None:
    manifest = store.ingest(pdf_path)
    store.save_annotation(
        manifest.material_id,
        Annotation(annotation_id="", locator=1, kind="note", note="whole-page note"),
    )
    store.save_annotation(
        manifest.material_id,
        Annotation(
            annotation_id="",
            locator=2,
            kind="highlight",
            quote="attention",
            rects=(Rect(0.2, 0.2, 0.7, 0.26),),
        ),
    )

    result = export_material(store, manifest.material_id, fmt="pdf")

    with pymupdf.open(stream=result.data, filetype="pdf") as doc:
        assert len(list(doc[0].annots())) == 1
        assert len(list(doc[1].annots())) == 1


def test_markdown_export_lists_marks_in_locator_order(store: ReadingStore, pdf_path: Path) -> None:
    manifest = store.ingest(pdf_path)
    store.save_annotation(
        manifest.material_id,
        Annotation(annotation_id="", locator=3, quote="Positional encoding", note="later"),
    )
    store.save_annotation(
        manifest.material_id,
        Annotation(annotation_id="", locator=1, quote="Introduction", note="first"),
    )

    result = export_material(store, manifest.material_id, fmt="markdown")
    text = result.data.decode("utf-8")

    assert result.filename == "attention-annotations.md"
    assert text.index("Page 1") < text.index("Page 3")
    assert "> Introduction" in text
    assert "first" in text


def test_markdown_export_labels_saved_citations(store: ReadingStore, pdf_path: Path) -> None:
    manifest = store.ingest(pdf_path)
    store.save_annotation(
        manifest.material_id,
        Annotation(
            annotation_id="",
            locator=1,
            kind="citation",
            quote="Introduction",
        ),
    )

    text = export_material(store, manifest.material_id, fmt="markdown").data.decode("utf-8")

    assert "**Citation**" in text
    assert "> Introduction" in text


def test_markdown_export_handles_a_material_with_no_annotations(
    store: ReadingStore, tmp_path: Path
) -> None:
    path = tmp_path / "notes.txt"
    path.write_text("some readable content here", encoding="utf-8")
    manifest = store.ingest(path)

    result = export_material(store, manifest.material_id, fmt="markdown")

    assert "No annotations yet" in result.data.decode("utf-8")


def test_auto_export_picks_pdf_for_pdfs_and_markdown_otherwise(
    store: ReadingStore, pdf_path: Path, tmp_path: Path
) -> None:
    pdf_manifest = store.ingest(pdf_path)
    text_path = tmp_path / "notes.txt"
    text_path.write_text("readable content", encoding="utf-8")
    text_manifest = store.ingest(text_path)

    assert export_material(store, pdf_manifest.material_id).media_type == "application/pdf"
    assert "markdown" in export_material(store, text_manifest.material_id).media_type


def test_pdf_export_is_refused_for_a_text_only_material(
    store: ReadingStore, tmp_path: Path
) -> None:
    path = tmp_path / "notes.txt"
    path.write_text("readable content", encoding="utf-8")
    manifest = store.ingest(path)

    with pytest.raises(ReadingError):
        export_material(store, manifest.material_id, fmt="pdf")


def test_state_written_before_the_split_is_still_read(store: ReadingStore, pdf_path: Path) -> None:
    """Annotations and viewports predate per-material state directories.

    They live beside the content as `annotations.json` / `position.json`, keyed
    by a material id that equalled the content hash. Those files must keep
    resolving, or every annotation a user made before the split disappears.
    """
    manifest = store.ingest(pdf_path)
    material_dir = store.root / manifest.material_id
    legacy_annotation = {
        "annotation_id": "legacy-1",
        "locator": 2,
        "quote": "scaled dot-product",
        "note": "written by the old reader",
        "color": "yellow",
        "kind": "highlight",
    }
    (material_dir / "annotations.json").write_text(
        json.dumps([legacy_annotation]), encoding="utf-8"
    )
    (material_dir / "position.json").write_text(
        json.dumps({"locator": 3, "source_anchor": "", "percentage": 0.0}),
        encoding="utf-8",
    )

    assert [row.note for row in store.annotations(manifest.material_id)] == [
        "written by the old reader"
    ]
    assert store.position(manifest.material_id).locator == 3

    # A write moves the material onto the per-material path without losing the
    # rows that were only in the legacy file.
    store.save_annotation(
        manifest.material_id,
        Annotation(
            annotation_id="new-1",
            locator=1,
            quote="sequence models",
            note="written after the split",
        ),
    )

    notes = sorted(row.note for row in store.annotations(manifest.material_id))
    assert notes == ["written after the split", "written by the old reader"]
