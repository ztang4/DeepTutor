"""Document text extraction for chat attachments.

Bytes-in, text-out. Used by the chat turn runtime to inline the text of
user-dropped files into the ``effective_user_message`` sent to the LLM.

Three format families:
  * **Binary Office** (.pdf / .docx / .xlsx / .pptx) — parsed with pymupdf /
    python-docx / openpyxl / python-pptx.
  * **EPUB** (.epub) — ZIP of XHTML documents; text is pulled in the OPF
    spine reading order using only the standard library.
  * **Text-like** (plain text, Markdown, source code, JSON, XML, CSV, …) —
    the extension set is imported from ``FileTypeRouter.TEXT_EXTENSIONS`` so
    the chat composer accepts every format the knowledge-base pipeline
    already ingests. Decoded with the same multi-encoding fallback chain.

Design mirrors ``nanobot/nanobot/utils/document.py`` but works on bytes
instead of file paths so the server never touches disk.
"""

from __future__ import annotations

import base64
from collections.abc import Iterable
from dataclasses import dataclass
from html.parser import HTMLParser
import io
import logging
from pathlib import Path, PurePosixPath
import posixpath
import re
from typing import Any
from urllib.parse import unquote
import zipfile

from defusedxml import ElementTree as DefusedElementTree
from defusedxml.common import DefusedXmlException

from deeptutor.services.rag.file_routing import FileTypeRouter

logger = logging.getLogger(__name__)

# Optional parser libraries are resolved on first use.  The public-ish module
# names remain overrideable because downstream deployments and tests use
# ``None`` to force the pure-OOXML fallback.
_NOT_LOADED = object()
fitz: Any = _NOT_LOADED
PdfReader: Any = _NOT_LOADED
_PypdfNotDecryptedError: Any = _NOT_LOADED
DocxDocument: Any = _NOT_LOADED
load_workbook: Any = _NOT_LOADED
PptxPresentation: Any = _NOT_LOADED


_OFFICE_EXTENSIONS: frozenset[str] = frozenset(FileTypeRouter.PARSER_EXTENSIONS)
# Text-like formats are sourced from the KB file router so chat and KB stay
# in sync. Adding a new code / config extension in one place propagates here.
TEXT_LIKE_EXTENSIONS: frozenset[str] = frozenset(FileTypeRouter.TEXT_EXTENSIONS)
SUPPORTED_DOC_EXTENSIONS: frozenset[str] = _OFFICE_EXTENSIONS | TEXT_LIKE_EXTENSIONS

# Built-in defaults, kept as module constants for callers that pass explicit
# budgets (the KB text_only engine, tests). The chat turn path resolves the
# effective values from system.json via ``_current_limits()`` on every call,
# so the /settings/attachments page applies without a restart.
MAX_DOC_BYTES = 20 * 1024 * 1024
MAX_TOTAL_DOC_BYTES = 25 * 1024 * 1024
MAX_EXTRACTED_CHARS_PER_DOC = 200_000
MAX_EXTRACTED_CHARS_TOTAL = 150_000


def _current_limits() -> tuple[int, int, int, int]:
    """(max_file_bytes, max_total_bytes, max_chars_per_doc, max_chars_total).

    Falls back to the module defaults when the settings layer is unavailable
    (e.g. unit tests running without a data directory).
    """
    try:
        from deeptutor.services.config.runtime_settings import get_chat_attachment_limits

        limits = get_chat_attachment_limits()
        return (
            limits.max_file_bytes,
            limits.max_total_bytes,
            limits.max_chars_per_doc,
            limits.max_chars_total,
        )
    except Exception:  # pragma: no cover - defensive fallback
        return (
            MAX_DOC_BYTES,
            MAX_TOTAL_DOC_BYTES,
            MAX_EXTRACTED_CHARS_PER_DOC,
            MAX_EXTRACTED_CHARS_TOTAL,
        )


_PDF_MAGIC = b"%PDF-"
_OOXML_MAGIC = b"PK\x03\x04"

_EPUB_CONTENT_EXTENSIONS: frozenset[str] = frozenset({".xhtml", ".html", ".htm"})
_EPUB_MAX_MEMBERS = 4096
_EPUB_MAX_MEMBER_BYTES = 20 * 1024 * 1024
_EPUB_MAX_TOTAL_UNCOMPRESSED_BYTES = 200 * 1024 * 1024
_EPUB_MAX_COMPRESSION_RATIO = 200.0
_EPUB_BLOCK_TAGS: frozenset[str] = frozenset(
    {
        "p",
        "div",
        "h1",
        "h2",
        "h3",
        "h4",
        "h5",
        "h6",
        "li",
        "blockquote",
        "tr",
        "table",
        "section",
        "article",
        "header",
        "footer",
        "aside",
        "figure",
        "figcaption",
        "dd",
        "dt",
        "dl",
        "hr",
    }
)


class DocumentExtractionError(Exception):
    """Base class for extraction failures. ``str(exc)`` is user-friendly."""

    def __init__(self, message: str, filename: str = "") -> None:
        super().__init__(message)
        self.filename = filename


class UnsupportedDocumentError(DocumentExtractionError):
    pass


class CorruptDocumentError(DocumentExtractionError):
    pass


class EmptyDocumentError(DocumentExtractionError):
    pass


class DocumentTooLargeError(DocumentExtractionError):
    pass


@dataclass(frozen=True, slots=True)
class EpubSpineUnit:
    """One EPUB spine document in package reading order.

    ``href`` is the normalised archive-member path. Keeping the source address
    next to its extracted text lets a faithful browser rendition and the
    server's numeric locator space refer to the same chapter.
    """

    href: str
    text: str
    title: str = ""


@dataclass(frozen=True, slots=True)
class EpubOutlineItem:
    """One EPUB navigation entry resolved to a spine locator."""

    locator: int
    title: str
    level: int = 1


def is_document_extension(filename: str) -> bool:
    return _ext(filename) in SUPPORTED_DOC_EXTENSIONS


def _ext(filename: str) -> str:
    return PurePosixPath(filename or "").suffix.lower()


def _truncate(text: str, max_length: int) -> str:
    if len(text) <= max_length:
        return text
    return text[:max_length] + f"... (truncated, {len(text)} chars total)"


def _check_magic(ext: str, data: bytes, filename: str) -> None:
    """Validate file header to catch extension spoofing.

    Only binary formats have well-known magic prefixes. Text-like extensions
    (code, markup, config, …) are decoded directly; a mislabeled binary blob
    either decodes as garbage or fails at decode time, which is fine.
    """
    if ext == ".pdf":
        if not data.startswith(_PDF_MAGIC):
            raise CorruptDocumentError(
                f"{filename} does not look like a PDF (bad header)", filename=filename
            )
    elif ext in {".docx", ".xlsx", ".pptx"}:
        if not data.startswith(_OOXML_MAGIC):
            raise CorruptDocumentError(
                f"{filename} does not look like a valid Office file (bad header)",
                filename=filename,
            )
    elif ext == ".epub":
        if not data.startswith(_OOXML_MAGIC):
            raise CorruptDocumentError(
                f"{filename} does not look like a valid EPUB (bad header)",
                filename=filename,
            )


def extract_text_from_bytes(
    filename: str,
    data: bytes,
    *,
    max_bytes: int | None = MAX_DOC_BYTES,
    max_chars: int | None = MAX_EXTRACTED_CHARS_PER_DOC,
) -> str:
    """Extract text from a single document's raw bytes.

    Raises a ``DocumentExtractionError`` subclass on failure. Successful
    output is truncated to ``max_chars`` with a notice when ``max_chars`` is
    not ``None``. ``max_bytes`` is configurable so the KB indexer can reuse
    the same parsers with its larger upload policy while chat keeps the
    stricter per-turn limit.
    """
    if not data:
        raise EmptyDocumentError(f"{filename} is empty", filename=filename)
    if max_bytes is not None and len(data) > max_bytes:
        raise DocumentTooLargeError(
            f"{filename} exceeds the {max_bytes // (1024 * 1024)} MB per-file limit",
            filename=filename,
        )

    ext = _ext(filename)
    if ext not in SUPPORTED_DOC_EXTENSIONS:
        raise UnsupportedDocumentError(
            f"{filename} has unsupported extension '{ext}'", filename=filename
        )

    _check_magic(ext, data, filename)

    if ext == ".pdf":
        text = _extract_pdf(data, filename)
    elif ext == ".docx":
        text = _extract_docx(data, filename)
    elif ext == ".xlsx":
        text = _extract_xlsx(data, filename)
    elif ext == ".pptx":
        text = _extract_pptx(data, filename)
    elif ext == ".epub":
        text = _extract_epub(data, filename)
    elif ext in TEXT_LIKE_EXTENSIONS:
        text = _extract_text_like(data, filename)
    else:  # pragma: no cover - guarded above
        raise UnsupportedDocumentError(f"{filename}: unreachable", filename=filename)

    if not text.strip():
        raise EmptyDocumentError(f"{filename}: no extractable text", filename=filename)
    return _truncate(text, max_chars) if max_chars is not None else text


def extract_text_from_path(
    file_path: str | Path,
    *,
    max_bytes: int | None = MAX_DOC_BYTES,
    max_chars: int | None = MAX_EXTRACTED_CHARS_PER_DOC,
) -> str:
    """Extract text from a file path using the same bytes-based parsers."""
    path = Path(file_path)
    return extract_text_from_bytes(
        path.name,
        path.read_bytes(),
        max_bytes=max_bytes,
        max_chars=max_chars,
    )


async def extract_text_from_path_isolated(
    file_path: str | Path,
    *,
    max_bytes: int | None = MAX_DOC_BYTES,
    max_chars: int | None = MAX_EXTRACTED_CHARS_PER_DOC,
    timeout: float = 120.0,
) -> str:
    """Extract in a short-lived spawn process and preserve public errors."""

    from deeptutor.runtime.isolated_worker import (
        IsolatedWorkerError,
        run_in_isolated_process,
    )

    path = Path(file_path)
    try:
        result = await run_in_isolated_process(
            "deeptutor.runtime.worker_tasks:extract_document_text",
            str(path),
            timeout=timeout,
            kwargs={"max_bytes": max_bytes, "max_chars": max_chars},
        )
    except IsolatedWorkerError as exc:
        error_types: dict[str, type[DocumentExtractionError]] = {
            cls.__name__: cls
            for cls in (
                DocumentExtractionError,
                UnsupportedDocumentError,
                CorruptDocumentError,
                EmptyDocumentError,
                DocumentTooLargeError,
            )
        }
        error_type = error_types.get(exc.remote_type)
        if error_type is not None:
            filename = str(exc.remote_attrs.get("filename") or path.name)
            raise error_type(str(exc), filename=filename) from exc
        if exc.remote_module == "builtins" and exc.remote_type in {
            "OSError",
            "FileNotFoundError",
            "PermissionError",
        }:
            raise OSError(str(exc)) from exc
        raise
    if not isinstance(result, str):
        raise DocumentExtractionError(
            f"{path.name}: isolated extractor returned invalid output",
            filename=path.name,
        )
    return result


def _extract_pdf(data: bytes, filename: str) -> str:
    global fitz, PdfReader, _PypdfNotDecryptedError
    if fitz is _NOT_LOADED:
        try:
            import fitz as fitz_module  # pymupdf

            fitz = fitz_module
        except ImportError:  # pragma: no cover
            fitz = None

    if fitz is not None:
        try:
            with fitz.open(stream=data, filetype="pdf") as doc:
                if doc.is_encrypted and not doc.authenticate(""):
                    raise CorruptDocumentError(
                        f"{filename} is encrypted and cannot be read", filename=filename
                    )
                pages = [
                    f"--- Page {i} ---\n{page.get_text() or ''}" for i, page in enumerate(doc, 1)
                ]
            return "\n\n".join(pages)
        except CorruptDocumentError:
            raise
        except Exception as exc:
            logger.warning("pymupdf failed on %s: %s — falling back to pypdf", filename, exc)

    if PdfReader is _NOT_LOADED:
        try:
            from pypdf import PdfReader as reader_type
            from pypdf.errors import FileNotDecryptedError

            PdfReader = reader_type
            _PypdfNotDecryptedError = FileNotDecryptedError
        except ImportError:  # pragma: no cover
            PdfReader = None
            _PypdfNotDecryptedError = Exception

    if PdfReader is None:
        raise CorruptDocumentError(
            f"{filename}: no PDF reader available (install pymupdf or pypdf)",
            filename=filename,
        )
    try:
        reader = PdfReader(io.BytesIO(data))
        if getattr(reader, "is_encrypted", False):
            raise CorruptDocumentError(
                f"{filename} is encrypted and cannot be read", filename=filename
            )
        pages = [
            f"--- Page {i} ---\n{page.extract_text() or ''}"
            for i, page in enumerate(reader.pages, 1)
        ]
        return "\n\n".join(pages)
    except CorruptDocumentError:
        raise
    except _PypdfNotDecryptedError as exc:
        raise CorruptDocumentError(
            f"{filename} is encrypted and cannot be read", filename=filename
        ) from exc
    except Exception as exc:
        raise CorruptDocumentError(
            f"{filename}: failed to read PDF ({exc})", filename=filename
        ) from exc


def _extract_docx(data: bytes, filename: str) -> str:
    global DocxDocument
    if DocxDocument is _NOT_LOADED:
        try:
            from docx import Document as document_type

            DocxDocument = document_type
        except ImportError:  # pragma: no cover
            DocxDocument = None

    primary_error: Exception | None = None
    primary_text = ""
    if DocxDocument is not None:
        try:
            doc = DocxDocument(io.BytesIO(data))
            paragraphs = [p.text for p in doc.paragraphs if p.text and p.text.strip()]
            primary_text = "\n\n".join(paragraphs)
        except Exception as exc:
            primary_error = exc
            logger.info("python-docx failed on %s; falling back to raw OOXML: %s", filename, exc)

    fallback = _extract_docx_ooxml(data, filename)
    if fallback.strip() and (not primary_text.strip() or len(fallback) > len(primary_text) * 1.2):
        return fallback
    if primary_text.strip():
        return primary_text

    if DocxDocument is None:
        raise CorruptDocumentError(
            f"{filename}: python-docx not installed and OOXML fallback found no text",
            filename=filename,
        )
    if primary_error is not None:
        raise CorruptDocumentError(
            f"{filename}: failed to open DOCX ({primary_error})", filename=filename
        ) from primary_error
    return ""


def _extract_xlsx(data: bytes, filename: str) -> str:
    global load_workbook
    if load_workbook is _NOT_LOADED:
        try:
            from openpyxl import load_workbook as workbook_loader

            load_workbook = workbook_loader
        except ImportError:  # pragma: no cover
            load_workbook = None

    if load_workbook is None:
        return _extract_xlsx_ooxml(data, filename)
    try:
        wb = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    except Exception as exc:
        logger.info("openpyxl failed on %s; falling back to raw OOXML: %s", filename, exc)
        fallback = _extract_xlsx_ooxml(data, filename)
        if fallback.strip():
            return fallback
        raise CorruptDocumentError(
            f"{filename}: failed to open XLSX ({exc})", filename=filename
        ) from exc
    try:
        sheets: list[str] = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            rows: list[str] = []
            for row in ws.iter_rows(values_only=True):
                row_text = "\t".join(str(cell) if cell is not None else "" for cell in row)
                if row_text.strip():
                    rows.append(row_text)
            if rows:
                sheets.append(f"--- Sheet: {sheet_name} ---\n" + "\n".join(rows))
        return "\n\n".join(sheets)
    finally:
        wb.close()


def _extract_pptx(data: bytes, filename: str) -> str:
    global PptxPresentation
    if PptxPresentation is _NOT_LOADED:
        try:
            from pptx import Presentation as presentation_type

            PptxPresentation = presentation_type
        except ImportError:  # pragma: no cover
            PptxPresentation = None

    if PptxPresentation is None:
        return _extract_pptx_ooxml(data, filename)
    try:
        prs = PptxPresentation(io.BytesIO(data))
    except Exception as exc:
        logger.info("python-pptx failed on %s; falling back to raw OOXML: %s", filename, exc)
        fallback = _extract_pptx_ooxml(data, filename)
        if fallback.strip():
            return fallback
        raise CorruptDocumentError(
            f"{filename}: failed to open PPTX ({exc})", filename=filename
        ) from exc
    slides: list[str] = []
    for i, slide in enumerate(prs.slides, 1):
        slide_text: list[str] = []
        for shape in slide.shapes:
            _collect_pptx_shape_text(shape, slide_text)
        if slide_text:
            slides.append(f"--- Slide {i} ---\n" + "\n".join(slide_text))
    return "\n\n".join(slides)


def _extract_text_like(data: bytes, filename: str) -> str:
    """Decode a plain-text / code / config / markup file.

    Uses the same encoding fallback chain as the KB pipeline
    (``FileTypeRouter.decode_bytes``) so a GBK-encoded Python file or a
    UTF-8-BOM Markdown works the same way in both places.
    """
    try:
        return FileTypeRouter.decode_bytes(data)
    except Exception as exc:  # pragma: no cover - decode_bytes never raises
        raise CorruptDocumentError(
            f"{filename}: failed to decode text ({exc})", filename=filename
        ) from exc


def _epub_parse_member(zf: zipfile.ZipFile, member: str, filename: str) -> Any | None:
    """Parse one XML/XHTML member, returning ``None`` when unreadable.

    Real-world EPUBs occasionally ship sloppy XHTML (undeclared entities,
    stray tags); a single bad chapter must not sink the whole book, so parse
    failures are logged and skipped instead of raising.
    """
    try:
        return _parse_xml_member(zf, member, filename)
    except CorruptDocumentError as exc:
        logger.warning("EPUB %s: skipping unparseable member %s (%s)", filename, member, exc)
        return None


class _EpubHTMLTextParser(HTMLParser):
    """Best-effort text renderer for EPUB chapters that are not valid XML."""

    def __init__(self) -> None:
        super().__init__(convert_charrefs=True)
        self.parts: list[str] = []
        self._ignored_depth = 0

    def handle_starttag(self, tag: str, attrs: list[tuple[str, str | None]]) -> None:
        del attrs
        tag = tag.lower()
        if tag in {"script", "style"}:
            self._ignored_depth += 1
        elif not self._ignored_depth and tag == "br":
            self.parts.append("\n")

    def handle_endtag(self, tag: str) -> None:
        tag = tag.lower()
        if tag in {"script", "style"} and self._ignored_depth:
            self._ignored_depth -= 1
        elif not self._ignored_depth and tag in _EPUB_BLOCK_TAGS:
            self.parts.append("\n\n")

    def handle_data(self, data: str) -> None:
        if not self._ignored_depth:
            self.parts.append(data)


def _epub_render_text(element: Any, parts: list[str]) -> None:
    """Append the text of one XHTML element and its subtree to ``parts``.

    Element text is emitted before its children and each child's tail after
    it, preserving the document's word spacing. Block-level tags contribute a
    paragraph break; ``script``/``style`` subtrees are dropped entirely.
    """
    tag = _local_name(element.tag) if isinstance(element.tag, str) else ""
    if tag in {"head", "script", "style"}:
        return
    if element.text:
        parts.append(element.text)
    if tag == "br":
        parts.append("\n")
    for child in element:
        _epub_render_text(child, parts)
        if child.tail:
            parts.append(child.tail)
    if tag in _EPUB_BLOCK_TAGS:
        parts.append("\n\n")


def _epub_xhtml_text(root: Any) -> str:
    """Render one XHTML document as plain text with paragraph breaks.

    Keeps the source whitespace of text nodes (XHTML carries its own word
    spacing) and collapses each line afterwards, so ``<b>world</b>.`` stays
    ``world.`` instead of gaining a stray space.
    """
    parts: list[str] = []
    _epub_render_text(root, parts)
    return _normalize_epub_text(parts)


def _normalize_epub_text(parts: Iterable[str]) -> str:
    raw = "".join(parts)
    lines = [re.sub(r"\s+", " ", line).strip() for line in raw.split("\n")]
    return "\n".join(line for line in lines if line)


def _epub_chapter_text(zf: zipfile.ZipFile, member: str, filename: str) -> str:
    """Render a chapter as XHTML, falling back to tolerant HTML parsing."""
    try:
        root = _parse_xml_member(zf, member, filename)
    except CorruptDocumentError as exc:
        try:
            raw = zf.read(member)
            parser = _EpubHTMLTextParser()
            parser.feed(FileTypeRouter.decode_bytes(raw))
            parser.close()
        except Exception:
            logger.warning("EPUB %s: skipping unparseable member %s", filename, member)
            return ""
        text = _normalize_epub_text(parser.parts)
        if text:
            logger.info("EPUB %s: used tolerant HTML parser for %s (%s)", filename, member, exc)
        return text
    return _epub_xhtml_text(root) if root is not None else ""


def _epub_html_members(names: list[str]) -> list[str]:
    """Archive members that look like XHTML content, in archive order."""
    return [name for name in names if _ext(name) in _EPUB_CONTENT_EXTENSIONS]


def _epub_content_files(zf: zipfile.ZipFile, filename: str) -> list[str]:
    """Resolve the XHTML content documents of an EPUB in reading order.

    Follows the standard chain ``META-INF/container.xml`` -> OPF package
    document -> spine ``itemref`` order. Falls back to every HTML/XHTML
    member in archive order when package metadata is missing or unusable.
    """
    names = zf.namelist()
    name_set = set(names)

    container_root = _epub_parse_member(zf, "META-INF/container.xml", filename)
    if container_root is None:
        return _epub_html_members(names)

    opf_path = ""
    for node in container_root.iter():
        if _local_name(node.tag) == "rootfile":
            opf_path = node.get("full-path") or ""
            break
    if not opf_path or opf_path not in name_set:
        return _epub_html_members(names)

    opf_root = _epub_parse_member(zf, opf_path, filename)
    if opf_root is None:
        return _epub_html_members(names)

    manifest: dict[str, str] = {}
    spine_ids: list[str] = []
    for node in opf_root.iter():
        name = _local_name(node.tag)
        if name == "item":
            item_id = node.get("id")
            href = node.get("href")
            if item_id and href:
                manifest[item_id] = href
        elif name == "itemref":
            idref = node.get("idref")
            if idref:
                spine_ids.append(idref)

    opf_dir = posixpath.dirname(opf_path)
    ordered: list[str] = []
    seen: set[str] = set()
    for idref in spine_ids:
        href = manifest.get(idref)
        if not href:
            continue
        member = posixpath.normpath(posixpath.join(opf_dir, unquote(href.split("#", 1)[0])))
        if member in name_set and member not in seen:
            ordered.append(member)
            seen.add(member)
    return ordered or _epub_html_members(names)


def _epub_element_text(node: Any) -> str:
    return re.sub(r"\s+", " ", "".join(node.itertext())).strip()


def _epub_package_navigation(
    zf: zipfile.ZipFile,
    filename: str,
    spine_members: list[str],
) -> list[EpubOutlineItem]:
    """Read EPUB3 nav or EPUB2 NCX entries and map them to spine locators."""
    container_root = _epub_parse_member(zf, "META-INF/container.xml", filename)
    if container_root is None:
        return []
    opf_path = next(
        (
            str(node.get("full-path") or "")
            for node in container_root.iter()
            if _local_name(node.tag) == "rootfile"
        ),
        "",
    )
    if not opf_path:
        return []
    opf_root = _epub_parse_member(zf, opf_path, filename)
    if opf_root is None:
        return []

    opf_dir = posixpath.dirname(opf_path)
    nav_member = ""
    ncx_member = ""
    spine_toc = ""
    manifest: dict[str, tuple[str, str]] = {}
    for node in opf_root.iter():
        name = _local_name(node.tag)
        if name == "item":
            item_id = str(node.get("id") or "")
            href = str(node.get("href") or "")
            properties = str(node.get("properties") or "")
            media_type = str(node.get("media-type") or "")
            if item_id and href:
                manifest[item_id] = (href, media_type)
                if "nav" in properties.split():
                    nav_member = posixpath.normpath(
                        posixpath.join(opf_dir, unquote(href.split("#", 1)[0]))
                    )
        elif name == "spine":
            spine_toc = str(node.get("toc") or "")

    if spine_toc in manifest:
        href, media_type = manifest[spine_toc]
        if media_type == "application/x-dtbncx+xml" or href.lower().endswith(".ncx"):
            ncx_member = posixpath.normpath(posixpath.join(opf_dir, unquote(href.split("#", 1)[0])))

    locator_by_member = {member: index for index, member in enumerate(spine_members, start=1)}

    def resolve_locator(href: str, base_member: str) -> int:
        path = unquote(href.split("#", 1)[0])
        member = posixpath.normpath(posixpath.join(posixpath.dirname(base_member), path))
        return locator_by_member.get(member, 0)

    rows: list[EpubOutlineItem] = []
    if nav_member:
        nav_root = _epub_parse_member(zf, nav_member, filename)
        if nav_root is not None:
            nav_nodes = [node for node in nav_root.iter() if _local_name(node.tag) == "nav"]
            toc_nav = next(
                (
                    node
                    for node in nav_nodes
                    if "toc"
                    in str(
                        node.get("{http://www.idpf.org/2007/ops}type")
                        or node.get("epub:type")
                        or ""
                    ).split()
                ),
                nav_nodes[0] if nav_nodes else None,
            )

            def walk_nav(node: Any, level: int) -> None:
                for child in node:
                    name = _local_name(child.tag)
                    if name == "li":
                        link = next(
                            (item for item in child if _local_name(item.tag) in {"a", "span"}),
                            None,
                        )
                        href = str(link.get("href") or "") if link is not None else ""
                        title = _epub_element_text(link) if link is not None else ""
                        locator = resolve_locator(href, nav_member) if href else 0
                        if locator and title:
                            rows.append(EpubOutlineItem(locator, title, max(1, level)))
                        for item in child:
                            if _local_name(item.tag) == "ol":
                                walk_nav(item, level + 1)
                    elif name == "ol":
                        walk_nav(child, level)

            if toc_nav is not None:
                walk_nav(toc_nav, 1)

    if not rows and ncx_member:
        ncx_root = _epub_parse_member(zf, ncx_member, filename)
        if ncx_root is not None:

            def walk_ncx(node: Any, level: int) -> None:
                for point in node:
                    if _local_name(point.tag) != "navPoint":
                        continue
                    label = next(
                        (item for item in point.iter() if _local_name(item.tag) == "navLabel"),
                        None,
                    )
                    content = next(
                        (item for item in point if _local_name(item.tag) == "content"),
                        None,
                    )
                    title = _epub_element_text(label) if label is not None else ""
                    href = str(content.get("src") or "") if content is not None else ""
                    locator = resolve_locator(href, ncx_member) if href else 0
                    if locator and title:
                        rows.append(EpubOutlineItem(locator, title, max(1, level)))
                    walk_ncx(point, level + 1)

            nav_map = next(
                (node for node in ncx_root.iter() if _local_name(node.tag) == "navMap"),
                None,
            )
            if nav_map is not None:
                walk_ncx(nav_map, 1)
    return rows


def extract_epub_spine(
    data: bytes,
    filename: str,
) -> tuple[tuple[EpubSpineUnit, ...], tuple[EpubOutlineItem, ...]]:
    """Return safe, source-addressed EPUB spine units and its nested outline."""
    _check_magic(".epub", data, filename)
    with _open_ooxml(data, filename) as zf:
        _validate_epub_archive(zf, filename)
        members = _epub_content_files(zf, filename)
        units: list[EpubSpineUnit] = []
        for member in members:
            text = _epub_chapter_text(zf, member, filename)
            heading = ""
            root = _epub_parse_member(zf, member, filename)
            if root is not None:
                heading_node = next(
                    (
                        node
                        for node in root.iter()
                        if _local_name(node.tag) in {"h1", "h2", "title"}
                        and _epub_element_text(node)
                    ),
                    None,
                )
                if heading_node is not None:
                    heading = _epub_element_text(heading_node)
            units.append(EpubSpineUnit(href=member, text=text, title=heading))
        outline = _epub_package_navigation(zf, filename, members)
    return tuple(units), tuple(outline)


def _extract_epub(data: bytes, filename: str) -> str:
    """Extract the reading text of an EPUB with only the standard library."""
    units, _ = extract_epub_spine(data, filename)
    return "\n\n".join(unit.text for unit in units if unit.text)


def _validate_epub_archive(zf: zipfile.ZipFile, filename: str) -> None:
    """Reject oversized or suspicious EPUB ZIPs before reading any member."""
    members = [info for info in zf.infolist() if not info.is_dir()]
    if len(members) > _EPUB_MAX_MEMBERS:
        raise DocumentTooLargeError(
            f"{filename}: EPUB has too many archive members ({len(members)})",
            filename=filename,
        )

    total = 0
    for info in members:
        if info.file_size > _EPUB_MAX_MEMBER_BYTES:
            raise DocumentTooLargeError(
                f"{filename}: EPUB member {info.filename} is too large",
                filename=filename,
            )
        total += info.file_size
        if total > _EPUB_MAX_TOTAL_UNCOMPRESSED_BYTES:
            raise DocumentTooLargeError(
                f"{filename}: EPUB uncompressed contents are too large",
                filename=filename,
            )
        if info.compress_size and info.file_size / info.compress_size > _EPUB_MAX_COMPRESSION_RATIO:
            raise DocumentTooLargeError(
                f"{filename}: EPUB member {info.filename} has a suspicious compression ratio",
                filename=filename,
            )


def _open_ooxml(data: bytes, filename: str) -> zipfile.ZipFile:
    try:
        return zipfile.ZipFile(io.BytesIO(data))
    except zipfile.BadZipFile as exc:
        raise CorruptDocumentError(
            f"{filename}: failed to open Office ZIP package ({exc})", filename=filename
        ) from exc


def _local_name(tag: str) -> str:
    return tag.rsplit("}", 1)[-1] if "}" in tag else tag


def _parse_xml_member(zf: zipfile.ZipFile, member: str, filename: str) -> Any | None:
    try:
        raw = zf.read(member)
    except KeyError:
        return None
    try:
        return DefusedElementTree.fromstring(raw)
    except (DefusedElementTree.ParseError, DefusedXmlException) as exc:
        raise CorruptDocumentError(
            f"{filename}: failed to parse {member} ({exc})", filename=filename
        ) from exc


def _collect_ooxml_text(node: Any) -> str:
    parts: list[str] = []
    for child in node.iter():
        name = _local_name(child.tag)
        if name == "t" and child.text:
            parts.append(child.text)
        elif name == "tab":
            parts.append("\t")
        elif name in {"br", "cr"}:
            parts.append("\n")
    return "".join(parts).strip()


def _extract_paragraph_text(root: Any) -> list[str]:
    paragraphs: list[str] = []
    for node in root.iter():
        if _local_name(node.tag) != "p":
            continue
        text = _collect_ooxml_text(node)
        if text:
            paragraphs.append(text)
    if paragraphs:
        return paragraphs
    text = _collect_ooxml_text(root)
    return [text] if text else []


def _extract_docx_ooxml(data: bytes, filename: str) -> str:
    with _open_ooxml(data, filename) as zf:
        names = zf.namelist()
        content_members = ["word/document.xml"]
        content_members.extend(
            sorted(
                name
                for name in names
                if re.match(r"word/(header|footer|footnotes|endnotes|comments)\d*\.xml$", name)
            )
        )

        chunks: list[str] = []
        for member in content_members:
            root = _parse_xml_member(zf, member, filename)
            if root is None:
                continue
            chunks.extend(_extract_paragraph_text(root))
        return "\n\n".join(chunks)


def _xlsx_shared_strings(zf: zipfile.ZipFile, filename: str) -> list[str]:
    root = _parse_xml_member(zf, "xl/sharedStrings.xml", filename)
    if root is None:
        return []
    strings: list[str] = []
    for node in root:
        if _local_name(node.tag) != "si":
            continue
        strings.append(_collect_ooxml_text(node))
    return strings


def _xlsx_sheet_names(zf: zipfile.ZipFile, filename: str) -> dict[str, str]:
    root = _parse_xml_member(zf, "xl/workbook.xml", filename)
    if root is None:
        return {}
    out: dict[str, str] = {}
    index = 1
    for node in root.iter():
        if _local_name(node.tag) != "sheet":
            continue
        sheet_name = node.attrib.get("name") or f"sheet{index}"
        sheet_id = node.attrib.get("sheetId") or str(index)
        out[f"xl/worksheets/sheet{sheet_id}.xml"] = sheet_name
        index += 1
    return out


def _xlsx_cell_text(cell: Any, shared_strings: list[str]) -> str:
    cell_type = cell.attrib.get("t", "")
    if cell_type == "inlineStr":
        return _collect_ooxml_text(cell)

    value = ""
    for child in cell:
        if _local_name(child.tag) == "v":
            value = child.text or ""
            break

    if cell_type == "s":
        try:
            return shared_strings[int(value)]
        except (ValueError, IndexError):
            return value
    return value


def _extract_xlsx_ooxml(data: bytes, filename: str) -> str:
    with _open_ooxml(data, filename) as zf:
        shared_strings = _xlsx_shared_strings(zf, filename)
        sheet_names = _xlsx_sheet_names(zf, filename)
        sheet_members = sorted(
            (name for name in zf.namelist() if re.match(r"xl/worksheets/sheet\d+\.xml$", name)),
            key=lambda name: [
                int(part) if part.isdigit() else part for part in re.split(r"(\d+)", name)
            ],
        )

        sheets: list[str] = []
        for index, member in enumerate(sheet_members, 1):
            root = _parse_xml_member(zf, member, filename)
            if root is None:
                continue
            rows: list[str] = []
            for row in root.iter():
                if _local_name(row.tag) != "row":
                    continue
                cells = [
                    _xlsx_cell_text(cell, shared_strings)
                    for cell in row
                    if _local_name(cell.tag) == "c"
                ]
                row_text = "\t".join(cells)
                if row_text.strip():
                    rows.append(row_text)
            if rows:
                sheet_name = sheet_names.get(member, f"sheet{index}")
                sheets.append(f"--- Sheet: {sheet_name} ---\n" + "\n".join(rows))
        return "\n\n".join(sheets)


def _extract_pptx_ooxml(data: bytes, filename: str) -> str:
    with _open_ooxml(data, filename) as zf:
        slide_members = sorted(
            (name for name in zf.namelist() if re.match(r"ppt/slides/slide\d+\.xml$", name)),
            key=lambda name: [
                int(part) if part.isdigit() else part for part in re.split(r"(\d+)", name)
            ],
        )
        slides: list[str] = []
        for index, member in enumerate(slide_members, 1):
            root = _parse_xml_member(zf, member, filename)
            if root is None:
                continue
            paragraphs = _extract_paragraph_text(root)
            if paragraphs:
                slides.append(f"--- Slide {index} ---\n" + "\n".join(paragraphs))
        return "\n\n".join(slides)


def _collect_pptx_shape_text(shape, out: list[str]) -> None:
    """Recurse into groups + tables, same semantics as nanobot's version."""
    sub_shapes = getattr(shape, "shapes", None)
    if sub_shapes is not None:
        for sub in sub_shapes:
            _collect_pptx_shape_text(sub, out)
        return

    if getattr(shape, "has_table", False):
        for row in shape.table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            line = "\t".join(cell for cell in cells if cell)
            if line:
                out.append(line)
        return

    text = getattr(shape, "text", "")
    if text:
        out.append(text)


def extract_documents_from_records(
    records: Iterable[dict],
) -> tuple[list[str], list[dict]]:
    """Process a list of attachment records from the WS payload.

    Parameters
    ----------
    records:
        Raw attachment records as parsed by the turn runtime
        (``{"type", "url", "base64", "filename", "mime_type"}``).

    Returns
    -------
    (doc_texts, updated_records)
        ``doc_texts`` is a list of strings formatted as
        ``"[File: <name>]\\n<text>"`` (one per processed or skipped doc).
        ``updated_records`` is the input list with the ``base64`` field
        cleared on successfully-extracted docs (to save DB space), an
        ``extracted_chars`` field added, and the extracted plain text
        stored under ``extracted_text`` so the chat UI can preview office
        documents without re-running the parser. Image / non-document
        records are returned unchanged.
    """
    doc_texts: list[str] = []
    updated: list[dict] = []
    max_file_bytes, max_total_bytes, max_chars_per_doc, max_chars_total = _current_limits()
    total_bytes = 0
    total_chars = 0
    over_quota = False

    for raw in records:
        record = dict(raw)
        filename = str(record.get("filename") or "")
        if not is_document_extension(filename):
            updated.append(record)
            continue

        b64 = record.get("base64") or ""
        if not b64:
            updated.append(record)
            continue

        if over_quota:
            doc_texts.append(f"[File: {filename} — skipped: total attachment quota exceeded]")
            record["base64"] = ""
            record["extracted_chars"] = 0
            updated.append(record)
            continue

        try:
            data = base64.b64decode(b64, validate=False)
        except Exception as exc:
            doc_texts.append(f"[File: {filename} — could not be read: invalid base64 ({exc})]")
            record["base64"] = ""
            record["extracted_chars"] = 0
            updated.append(record)
            continue

        if total_bytes + len(data) > max_total_bytes:
            over_quota = True
            doc_texts.append(f"[File: {filename} — skipped: total attachment quota exceeded]")
            record["base64"] = ""
            record["extracted_chars"] = 0
            updated.append(record)
            continue

        total_bytes += len(data)

        try:
            text = extract_text_from_bytes(
                filename,
                data,
                max_bytes=max_file_bytes,
                max_chars=max_chars_per_doc,
            )
        except DocumentExtractionError as exc:
            logger.info("Document extraction failed for %s: %s", filename, exc)
            doc_texts.append(f"[File: {filename} — could not be read: {exc}]")
            record["base64"] = ""
            record["extracted_chars"] = 0
            updated.append(record)
            continue

        remaining_budget = max_chars_total - total_chars
        if remaining_budget <= 0:
            doc_texts.append(f"[File: {filename} — skipped: total extracted-text quota exceeded]")
            record["base64"] = ""
            record["extracted_chars"] = 0
            updated.append(record)
            continue

        if len(text) > remaining_budget:
            text = (
                text[:remaining_budget]
                + f"... (truncated, {len(text)} chars total; turn quota hit)"
            )

        total_chars += len(text)
        doc_texts.append(f"[File: {filename}]\n{text}")
        record["base64"] = ""
        record["extracted_chars"] = len(text)
        record["extracted_text"] = text
        updated.append(record)

    return doc_texts, updated
